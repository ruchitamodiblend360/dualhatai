"""
Local dev server for the User Story Readiness Checker.

- Serves index.html (the React UI) at /
- Proxies POST /api/analyze to an LLM (OpenAI by default, Groq as automatic
  fallback), injecting the API keys from .env server-side so they never touch
  the browser.

Run:  py server.py    (or: python server.py)
Then open http://localhost:8000  in your browser.
"""

import json
import os
import re
import time
import urllib.request
import urllib.error
from collections import Counter
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from concurrent.futures import ThreadPoolExecutor, as_completed
import datetime
from zoneinfo import ZoneInfo

EASTERN_TZ = ZoneInfo("America/New_York")

ROOT = os.path.dirname(os.path.abspath(__file__))
HISTORY_FILE = os.path.join(ROOT, "history.json")


def load_history():
    if not os.path.exists(HISTORY_FILE):
        return []
    try:
        with open(HISTORY_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []


def save_history(entries):
    with open(HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(entries, f, ensure_ascii=False, indent=2)


_STOP_WORDS = {
    "a","an","the","as","is","in","of","to","and","or","for","with","that",
    "this","we","i","it","be","at","by","on","so","our","can","are","have",
    "has","from","their","will","not","should","when","was","were","if","but",
    "about","into","they","he","she","its","all","any","each","get","use",
    "new","my","your","his","her","what","which","who","do","how","may","more",
    "also","want","user","story","epic","able","would","could","need","than",
}

def _keyword_set(text):
    return {w for w in re.findall(r'\b[a-z]{3,}\b', text.lower()) if w not in _STOP_WORDS}

def _jaccard(a, b):
    sa, sb = _keyword_set(a), _keyword_set(b)
    if not sa or not sb:
        return 0.0
    return len(sa & sb) / len(sa | sb)


def build_team_context(history, mode, current_story, max_examples=2):
    """Build Level 1 (few-shot examples) + Level 2 (pattern summary) context block."""
    seen_texts = set()
    relevant = []
    for e in history:
        if not (e.get("result") and (e.get("mode") or "story") == mode and e.get("story_text")):
            continue
        key = e.get("story_text", "")[:120]
        if key in seen_texts:
            continue
        seen_texts.add(key)
        relevant.append(e)
    if not relevant:
        return ""

    parts = []
    mode_label = "stories" if mode == "story" else "epics"

    # ── Level 2: pattern summary (requires 3+ past entries) ──────────────────
    if len(relevant) >= 3:
        dim_names = ["completeness", "clarity", "testability", "size", "dependency_risk"]
        all_scores = [e["result"].get("scores", {}) for e in relevant]
        totals     = [e["result"]["total"] for e in relevant if e["result"].get("total")]

        dim_avgs = {}
        for d in dim_names:
            vals = [s[d] for s in all_scores if isinstance(s.get(d), (int, float))]
            if vals:
                dim_avgs[d] = round(sum(vals) / len(vals), 1)

        weak_dims = [
            d.replace("_", " ")
            for d, v in sorted(dim_avgs.items(), key=lambda x: x[1])
            if v < 12
        ]

        gap_areas = [
            g["area"]
            for e in relevant
            for g in (e["result"].get("gaps") or [])
            if g.get("area")
        ]
        top_gap_areas = [area for area, _ in Counter(gap_areas).most_common(3)]

        avg_total = round(sum(totals) / len(totals), 1) if totals else None

        lines = [f"TEAM HISTORICAL PATTERNS ({len(relevant)} past {mode_label} analysed):"]
        if avg_total is not None:
            lines.append(f"- Team average score: {avg_total}/100")
        if weak_dims:
            lines.append(
                f"- Consistently weak dimensions: {', '.join(weak_dims)}"
                f" - apply stricter scrutiny here"
            )
        if top_gap_areas:
            lines.append(f"- Most recurring gap areas: {', '.join(top_gap_areas)}")
        lines.append(
            "Calibrate your scoring against these patterns. "
            "Be especially critical in the weak dimensions listed above."
        )
        parts.append("\n".join(lines))

    # ── Level 1: few-shot similar examples ───────────────────────────────────
    ranked = sorted(relevant, key=lambda e: -_jaccard(current_story, e["story_text"]))
    examples = [e for e in ranked[:max_examples] if _jaccard(current_story, e["story_text"]) > 0.05]

    if examples:
        ex_lines = [f"SIMILAR PAST {mode_label.upper()} FROM THIS TEAM (score calibration):"]
        for e in examples:
            r = e["result"]
            preview = (e["story_text"] or "")[:200].replace("\n", " ").strip()
            if len(e["story_text"]) > 200:
                preview += "…"
            gap_summary = "; ".join(
                f"{g['area']}: {g['issue'][:70]}"
                for g in (r.get("gaps") or [])[:2]
                if g.get("area") and g.get("issue")
            ) or "none noted"
            ex_lines.append(
                f'Story: "{preview}"\n'
                f'Score: {r.get("total","?")} / 100 | Level: {r.get("readiness_level","?")}\n'
                f'Key gaps: {gap_summary}'
            )
        parts.append("\n\n".join(ex_lines))

    return "\n\n---\n\n".join(parts) if parts else ""


GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "llama-3.3-70b-versatile"
MODEL = GROQ_MODEL  # backward-compat alias
OPENAI_URL = "https://api.openai.com/v1/chat/completions"
OPENAI_MODEL = "gpt-4o-mini"

# ── User Story prompt ────────────────────────────────────────────────────────
SYSTEM_PROMPT = """You are an expert agile coach and product manager with 15+ years of experience running sprint planning, backlog refinement, and Definition of Ready reviews across engineering teams.

Your job is to analyse a user story submitted before sprint planning and return a structured readiness report. You evaluate stories the way a senior agile practitioner would — not just checking format, but assessing whether the story gives a development team everything they need to build, test, and ship without ambiguity or mid-sprint blockers.

TERMINOLOGY: The submission you are analysing is a USER STORY, not an epic. In every free-text field you write (summary, gaps, ambiguities, dependencies, improved_story), always refer to it as "the story" or "the user story" — never as "the epic".

---

## SCORING DIMENSIONS

Score the story across exactly 5 dimensions. Each dimension is scored 0–20. Total score is 0–100.

### 1. COMPLETENESS (0–20)
What to check:
- Is the "As a [persona]… I want [goal]… so that [benefit]" format present and meaningful?
- Is the persona specific (not generic like "user" or "admin")?
- Are acceptance criteria present? Do they cover the main flow?
- Are edge cases and error states defined?
- Is out-of-scope explicitly stated where needed?

Scoring guide:
- 0–5:   No format, no ACs, no goal defined
- 6–10:  Format partially present, ACs missing or very thin
- 11–15: Good structure, some edge cases missing
- 16–20: Full format, ACs defined, major flows covered. Edge cases are nice but not required for sprint-readiness.

### 2. CLARITY (0–20)
What to check:
- Is any language vague or subjective without a measurable definition?
- Watch for: "fast", "quick", "easy", "intuitive", "user-friendly", "simple", "better", "improved", "seamless", "robust", "scalable", "secure" — all flagged unless defined with a metric
- Are all terms consistently used? Are abbreviations or system names explained?
- Can a developer who is new to the team understand this story without asking questions?

Scoring guide:
- 0–5:   Multiple undefined subjective terms, story cannot be understood without follow-up
- 6–10:  Several vague terms, likely to cause mid-sprint clarification requests
- 11–15: Minor unclear terms, mostly understandable
- 16–20: Precise language throughout, all terms measurable or clearly defined

### 3. TESTABILITY (0–20)
What to check:
- Can every acceptance criterion be independently verified by a QA engineer?
- Are ACs written in Given/When/Then (or equivalent) format?
- Are there measurable outcomes (SLAs, counts, specific error messages, HTTP status codes)?
- Are success and failure paths both defined?
- Could a test be written for each AC without asking the author for clarification?

Scoring guide:
- 0–5:   ACs are absent or entirely untestable ("system works correctly")
- 6–10:  ACs present but mostly unverifiable without further information
- 11–15: Most ACs testable, some gaps in failure paths or edge cases
- 16–20: All ACs independently verifiable, Given/When/Then, measurable outcomes throughout

### 4. SIZE (0–20)
What to check:
- Does the story represent a single, deliverable unit of value?
- Can it realistically be completed in one sprint by one team?
- Is it an epic in disguise (covering multiple features or flows)?
- Could it be split into smaller independently deliverable stories?
- Is a story point estimate included or inferable?

Scoring guide:
- 0–5:   Clearly an epic in disguise — covers 3+ distinct features or workflows
- 6–10:  Too large for a single sprint, should be split
- 11–15: Borderline — could be completed in a sprint but is on the larger side
- 16–20: Well-scoped, single unit of value, sprint-sized

### 5. DEPENDENCY RISK (0–20)
What to check:
- Are there implied or explicit dependencies on other teams, APIs, services, or stories?
- Are named dependencies acknowledged with owning team or contact?
- Are there integration points that could block progress mid-sprint?
- Are there design, legal, compliance, or infrastructure dependencies?
- Are dependencies de-risked (e.g. contract agreed, API documented, designs approved)?

Scoring guide:
- 0–5:   Multiple unacknowledged dependencies that are likely blockers
- 6–10:  Dependencies implied but not named; risk of mid-sprint blocks is high
- 11–15: Some dependencies named; risk is medium; follow-up needed
- 16–20: Dependencies fully acknowledged, named, and de-risked or low-risk

---

## OUTPUT FORMAT

Return ONLY a single valid JSON object. No markdown. No code fences. No explanation before or after. No comments inside the JSON.

The JSON must exactly match this structure:

{
  "scores": {
    "completeness": <integer 0–20>,
    "clarity": <integer 0–20>,
    "testability": <integer 0–20>,
    "size": <integer 0–20>,
    "dependency_risk": <integer 0–20>
  },
  "total": <integer 0–100, must equal sum of all 5 scores>,
  "readiness_level": "<exactly one of: Not Ready | Needs Work | Almost Ready | Sprint Ready>",
  "summary": "<2–3 sentences. State the overall verdict, the 1–2 biggest strengths, and the 1–2 most critical issues. Refer to the submission as "the story". Be direct and specific — avoid generic statements.>",
  "gaps": [
    {
      "severity": "<exactly one of: critical | warning | info>",
      "area": "<which dimension this gap belongs to: Completeness | Clarity | Testability | Size | Dependencies>",
      "issue": "<specific description of the gap — quote the problematic phrase or missing element>",
      "fix": "<concrete, actionable fix — tell the team exactly what to add, change, or define>"
    }
  ],
  "ambiguities": [
    {
      "phrase": "<exact phrase from the story that is ambiguous>",
      "question": "<the specific question the team must answer before this story is sprint-ready>"
    }
  ],
  "dependencies": [
    {
      "type": "<exactly one of: team | api | story | system | design | compliance>",
      "description": "<what the dependency is and why it could block progress>",
      "confidence": "<exactly one of: high | medium | low>",
      "status": "<exactly one of: acknowledged | implied | unresolved>"
    }
  ],
  "improved_story": "<Full rewrite of the story in correct format. Preserve the original intent. Fix vague language with measurable alternatives. Do not add acceptance criteria here — those go in suggested_acs.>",
  "suggested_acs": [
    "<AC in Given [context] / When [action] / Then [measurable outcome] format>",
    "<include at least 3, up to 7 ACs covering happy path, error states, and edge cases>"
  ]
}

---

## READINESS LEVEL THRESHOLDS

Map total score to readiness_level as follows:
- 0–39:   "Not Ready"
- 40–59:  "Needs Work"
- 60–79:  "Almost Ready"
- 80–100: "Sprint Ready"

---

## SEVERITY DEFINITIONS

Use these definitions consistently across all gap entries:

- critical: A blocker. The story cannot be built or tested without resolving this. Examples: missing ACs, undefined token expiry, no error states, epic-sized scope.
- warning:  Likely to cause a mid-sprint clarification request or rework. Examples: vague language with no metric, missing edge case, dependency named but not de-risked.
- info:     Nice to have. Will improve quality but won't block delivery. Examples: out-of-scope not stated, story point estimate missing, minor inconsistency in terminology.

---

## TEAM CONTEXT (if provided)

If the user provides team context (Definition of Ready, parent epic, story point scale, team conventions), incorporate it into your scoring:
- Score the story against the team's stated DoR, not a generic one
- Flag any DoR criteria the story fails to meet as critical gaps
- Reference the parent epic (if provided) when assessing dependency risk and scope

---

## BEHAVIOUR RULES

1. Never invent information. If something is not in the story, flag it as missing — do not assume it exists.
2. Be specific. Quote exact phrases when flagging ambiguities or gaps. "User can log in quickly" is a quote; "vague language present" is not useful.
3. Be proportionate. A story with 1 minor vague word should not score the same as a story with no ACs at all.
4. Gaps array: include only gaps that would cause mid-sprint blockers or create ambiguity for the engineering team. A missing edge case detail is not a gap; a missing acceptance criteria is. Maximum 5 gaps.
5. Ambiguities array: only include phrases that are genuinely ambiguous. Do not manufacture ambiguity in an otherwise clear story.
6. Dependencies array: include both explicit (named in the story) and strongly implied dependencies. Set confidence to "low" for implied ones.
7. Improved story: rewrite the story narrative only. Do not insert ACs into the improved_story field — they belong in suggested_acs.
8. Total score must arithmetically equal the sum of the five dimension scores.
9. Terminology: never call the submission "the epic" in any free-text field — it is "the story" or "the user story".
10. Jira comments (if provided): treat clarifying questions, scope disagreements, or dependencies raised in them as real evidence for gaps, ambiguities, or dependencies — a story that reads cleanly but generated confused comments is not actually clear. Quote the comment the same way you'd quote the story text.
11. Return valid JSON only. Any deviation breaks the application."""

# ── Epic prompt ──────────────────────────────────────────────────────────────
EPIC_SYSTEM_PROMPT = """You are an expert agile coach, product manager, and software architect with 15+ years of experience running sprint planning, backlog refinement, and Definition of Ready reviews across engineering teams.

Your job is to analyse an epic submitted before sprint planning and return a structured readiness report. You evaluate epics the way a senior agile practitioner would — not just checking format, but assessing whether the epic gives a development team everything they need to build, test, and ship without ambiguity or mid-sprint blockers.

TERMINOLOGY: The submission you are analysing is an EPIC, not a user story. In every free-text field you write (summary, gaps, ambiguities, dependencies, improved_story, split_suggestions), always refer to it as "the epic" — never as "the story" or "the user story".

---

## SCORING DIMENSIONS

Score the epic across exactly 5 dimensions. Each dimension is scored 0–20. Total score is 0–100.

### 1. COMPLETENESS (0–20)
What to check:
- Is the epic's goal and business outcome clearly stated (not just a feature name)?
- Is the scope defined — what's included, and what's explicitly out of scope?
- Are the key stakeholders or impacted teams identified?
- Is the format meaningful (not just a title with no context)?

Scoring guide:
- 0–5:   No goal, outcome, or scope stated
- 6–10:  Goal stated but business outcome or scope is thin/unclear
- 11–15: Goal, outcome, and scope present; stakeholders or out-of-scope not stated
- 16–20: Full goal, business outcome, scope, and out-of-scope all stated. Granular edge cases and Definition of Done are scored under TESTABILITY, not here.

### 2. CLARITY (0–20)
What to check:
- Is any language vague or subjective without a measurable definition?
- Watch for: "fast", "quick", "easy", "intuitive", "user-friendly", "simple", "better", "improved", "seamless", "robust", "scalable", "secure" — all flagged unless defined with a metric
- Are all terms consistently used? Are abbreviations or system names explained?
- Can a developer who is new to the team understand this epic without asking questions?

Scoring guide:
- 0–5:   Multiple undefined subjective terms, epic cannot be understood without follow-up
- 6–10:  Several vague terms, likely to cause mid-sprint clarification requests
- 11–15: Minor unclear terms, mostly understandable
- 16–20: Precise language throughout, all terms measurable or clearly defined

### 3. TESTABILITY (0–20)
An epic itself rarely has granular acceptance criteria — what matters is whether the team could tell, unambiguously, when the epic as a whole is done, and whether its goals are concrete enough for testable child stories to be written from them.

What to check:
- Is there a clear, epic-level Definition of Done — a set of measurable conditions that define when the epic is complete (not just "when all child stories are done", but what those stories collectively need to prove)?
- Is the "done" state observable and verifiable (specific metrics, launched capabilities, a named set of child stories or milestones) rather than subjective ("users are happier", "the system is better")?
- Are the epic's stated goals concrete enough that a team could write independently testable acceptance criteria for each child story without inventing scope?
- If child stories or milestones are listed or implied, is there enough detail to judge whether each would be testable on its own?

Scoring guide:
- 0–5:   No definition of done; success criteria are subjective or entirely absent
- 6–10:  A rough goal is stated but not measurable; unclear how anyone would verify the epic is complete
- 11–15: Definition of done is mostly measurable, some gaps in how child-story-level testability would work
- 16–20: Clear, measurable epic-level definition of done; goals are concrete enough that testable child stories could be written directly from them

### 4. SIZE (0–20)
An epic is expected to span multiple sprints — do not penalize it for being larger than a single story. What matters is whether its scope is bounded and cleanly decomposable, not whether it fits in one sprint.

What to check:
- Is the epic bounded — does it have a clear start and end state, or is it open-ended ("and more", "etc.", unscoped future work folded in)?
- Does it represent a single coherent initiative or theme, rather than several unrelated initiatives bundled together under one epic?
- Is there a rough scope signal — an estimated number of sprints, a story count, or a team-months estimate — stated or clearly inferable?
- Could it realistically be decomposed into a set of sprint-sized child stories right now, without requiring another round of scoping first?

Scoring guide:
- 0–5:   Unbounded scope with no defined end state, or bundles multiple unrelated initiatives together
- 6–10:  Bounded but too coarse to decompose yet — needs another scoping pass before child stories can be written
- 11–15: Mostly bounded and decomposable, but a scope estimate is missing or the boundary is fuzzy in places
- 16–20: Clearly bounded single initiative, scope estimate present, cleanly decomposable into sprint-sized stories today

### 5. DEPENDENCY RISK (0–20)
What to check:
- Are there implied or explicit dependencies on other teams, APIs, services, or stories?
- Are named dependencies acknowledged with owning team or contact?
- Are there integration points that could block progress mid-sprint?
- Are there design, legal, compliance, or infrastructure dependencies?
- Are dependencies de-risked (e.g. contract agreed, API documented, designs approved)?

Scoring guide:
- 0–5:   Multiple unacknowledged dependencies that are likely blockers
- 6–10:  Dependencies implied but not named; risk of mid-sprint blocks is high
- 11–15: Some dependencies named; risk is medium; follow-up needed
- 16–20: Dependencies fully acknowledged, named, and de-risked or low-risk

---

## OUTPUT FORMAT

Return ONLY a single valid JSON object. No markdown. No code fences. No explanation before or after. No comments inside the JSON.

The JSON must exactly match this structure:

{
  "scores": {
    "completeness": <integer 0–20>,
    "clarity": <integer 0–20>,
    "testability": <integer 0–20>,
    "size": <integer 0–20>,
    "dependency_risk": <integer 0–20>
  },
  "total": <integer 0–100, must equal sum of all 5 scores>,
  "readiness_level": "<exactly one of: Not Ready | Needs Work | Almost Ready | Decomposition Ready>",
  "summary": "<2–3 sentences. State the overall verdict, the 1–2 biggest strengths, and the 1–2 most critical issues. Refer to the submission as "the epic". Be direct and specific — avoid generic statements.>",
  "gaps": [
    {
      "severity": "<exactly one of: critical | warning | info>",
      "area": "<which dimension this gap belongs to: Completeness | Clarity | Testability | Size | Dependencies>",
      "issue": "<specific description of the gap — quote the problematic phrase or missing element>",
      "fix": "<concrete, actionable fix — tell the team exactly what to add, change, or define>"
    }
  ],
  "ambiguities": [
    {
      "phrase": "<exact phrase from the epic that is ambiguous>",
      "question": "<the specific question the team must answer before this epic is ready to decompose into stories>"
    }
  ],
  "dependencies": [
    {
      "type": "<exactly one of: team | api | story | system | design | compliance>",
      "description": "<what the dependency is and why it could block progress>",
      "confidence": "<exactly one of: high | medium | low>",
      "status": "<exactly one of: acknowledged | implied | unresolved>"
    }
  ],
  "improved_story": "<Full rewrite of the epic in correct format. Preserve the original intent. Fix vague language with measurable alternatives. Do not add acceptance criteria here — those go in suggested_acs.>",
  "suggested_acs": [
    "<NOT granular feature-level ACs — an epic doesn't need those. Instead, write epic-level Definition of Done criteria: the measurable conditions that prove the epic as a whole is complete (e.g. 'Given all committed capabilities have shipped, when [metric/capability] is verified, then the epic is done'). 3–5 items, each stated at the epic level, consistent with the TESTABILITY dimension above.>"
  ],
  "split_suggestions": [
    "<Always populate this for an epic, regardless of the size score — decomposing into child stories is expected output, not a fallback for a poorly-scoped one. Suggest 3–6 child stories the epic could be decomposed into, each a one-sentence story title. If size score is ≤ 10, first note the missing boundary/estimate as a gap, then still provide the best decomposition you can from what's given.>"
  ]
}

---

## READINESS LEVEL THRESHOLDS

Map total score to readiness_level as follows:
- 0–39:   "Not Ready"
- 40–59:  "Needs Work"
- 60–79:  "Almost Ready"
- 80–100: "Decomposition Ready"

---

## SEVERITY DEFINITIONS

Use these definitions consistently across all gap entries:

- critical: A blocker. The epic cannot be decomposed or built without resolving this. Examples: missing ACs, undefined token expiry, no error states, unbounded scope.
- warning:  Likely to cause a mid-sprint clarification request or rework. Examples: vague language with no metric, missing edge case, dependency named but not de-risked.
- info:     Nice to have. Will improve quality but won't block delivery. Examples: out-of-scope not stated, story point estimate missing, minor inconsistency in terminology.

---

## TEAM CONTEXT (if provided)

If the user provides team context (Definition of Ready, parent epic, story point scale, team conventions), incorporate it into your scoring:
- Score the epic against the team's stated DoR, not a generic one
- Flag any DoR criteria the epic fails to meet as critical gaps
- Reference the parent epic context (if provided) when assessing dependency risk and scope

---

## BEHAVIOUR RULES

1. Never invent information. If something is not in the epic, flag it as missing — do not assume it exists.
2. Be specific. Quote exact phrases when flagging ambiguities or gaps. "User can log in quickly" is a quote; "vague language present" is not useful.
3. Be proportionate. An epic with 1 minor vague word should not score the same as an epic with no Definition of Done at all.
4. Gaps array: include only gaps that would cause mid-sprint blockers or create ambiguity for the engineering team. A missing edge case detail is not a gap; a missing acceptance criteria is. Maximum 5 gaps.
5. Ambiguities array: only include phrases that are genuinely ambiguous. Do not manufacture ambiguity in an otherwise clear epic.
6. Dependencies array: include both explicit (named in the epic) and strongly implied dependencies. Set confidence to "low" for implied ones.
7. Improved epic: rewrite the epic narrative only. Do not insert Definition of Done criteria into the improved_story field — they belong in suggested_acs.
8. Split suggestions: always populate with a decomposition into child stories — do not gate this on the size score.
9. Total score must arithmetically equal the sum of the five dimension scores.
10. Terminology: never call the submission "the story" or "the user story" in any free-text field — it is "the epic".
11. Jira comments (if provided): treat clarifying questions, scope disagreements, or dependencies raised in them as real evidence for gaps, ambiguities, or dependencies — an epic that reads cleanly but generated confused comments is not actually clear. Quote the comment the same way you'd quote the epic text.
12. Return valid JSON only. Any deviation breaks the application."""


# ── Status Deck prompt ──────────────────────────────────────────────────────
STATUS_DECK_SYSTEM_PROMPT = """You are an executive project status generator. Given Jira sprint issue data, synthesize a concise, stakeholder-ready weekly status deck.

Output rules:

1. executiveSummary: 2-3 sentences. State overall project health, the primary focus/win this week, and any critical concern or blocker.

2. accomplishments: Extract from Done issues. For each item the "sentence" field MUST be a past-tense sentence YOU COMPOSE (8-15 words) starting with a strong action verb — do NOT copy the Jira summary verbatim. Use the "Description:" text to build a specific, meaningful sentence about what was actually delivered. If no description, craft the best sentence from the summary.
   Strong verbs: Implemented, Deployed, Resolved, Delivered, Launched, Fixed, Migrated, Integrated, Optimized, Shipped, Automated, Configured, Built, Enabled, Completed, Released, Refactored.
   Examples: summary 'User auth API' + desc 'JWT auth with refresh tokens' -> sentence 'Implemented JWT-based user authentication with automatic token refresh'. Summary 'Fix login bug' + desc 'Fixed redirect loop' -> sentence 'Resolved OAuth callback redirect loop causing login failures'. 4-8 items max.

3. nextSteps: Extract from In Progress + high-priority To Do. Action-oriented, starting with a verb. Include owner and due date where derivable — if the ticket fields don't specify one, check "Latest comments" for a stated next action or owner before leaving it blank. 4-8 items max. Never generic ("continue development").

4. blockers: Items labeled blocked or risk, unassigned critical items, issues with no progress. Also include issues whose "Latest comments" reveal they're stuck, waiting on another team, or have an unresolved question — even if status/labels don't say so. When comments name who owns the fix or describe a resolution plan, use that for "owner" and "mitigation" instead of leaving them blank. 5 max.

5. healthStatus: On Track if >60% complete and no critical blockers. At Risk if 40-60% or has blockers. Off Track if <40% or multiple critical blockers.

6. milestones: Extract sprint goals or milestone-labeled issues. Return [] if none.

Return ONLY valid JSON, no markdown:
{
  "projectName": "string",
  "sprintName": "string",
  "weekOf": "string",
  "healthStatus": "On Track | At Risk | Off Track",
  "executiveSummary": "string",
  "accomplishments": [{"sentence": "string (past-tense, AI-composed, NOT the raw Jira summary)"}],
  "nextSteps": [{"action": "string", "owner": "string", "dueDate": "string", "priority": "high|medium|low"}],
  "blockers": [{"title": "string", "impact": "High|Medium|Low", "type": "blocker|risk|dependency", "mitigation": "string", "owner": "string"}],
  "milestones": [{"name": "string", "status": "complete|in_progress|upcoming", "date": "string"}]
}"""


def load_env():
    path = os.path.join(ROOT, ".env")
    cfg = {}
    if not os.path.exists(path):
        return cfg
    with open(path, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            k, v = line.split("=", 1)
            cfg[k.strip()] = v.strip().strip('"').strip("'")
    return cfg


_ENV = load_env()
GROQ_API_KEY   = _ENV.get("GROQ_API_KEY")
OPENAI_API_KEY = _ENV.get("OPENAI_API_KEY")
if _ENV.get("OPENAI_MODEL"):
    OPENAI_MODEL = _ENV.get("OPENAI_MODEL")
if _ENV.get("GROQ_MODEL"):
    GROQ_MODEL = _ENV.get("GROQ_MODEL")
JIRA_BASE_URL  = _ENV.get("JIRA_BASE_URL", "").rstrip("/")
JIRA_EMAIL     = _ENV.get("JIRA_EMAIL", "")
JIRA_API_TOKEN = _ENV.get("JIRA_API_TOKEN", "")

import base64 as _b64
JIRA_AUTH = _b64.b64encode(f"{JIRA_EMAIL}:{JIRA_API_TOKEN}".encode()).decode() if JIRA_EMAIL and JIRA_API_TOKEN else None

LLM_UA = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) story-readiness-checker/1.0"


class LLMError(Exception):
    def __init__(self, status, detail):
        super().__init__(detail)
        self.status = status
        self.detail = detail


def _llm_post(url, api_key, body_dict, timeout=90):
    data = json.dumps(body_dict).encode("utf-8")
    req = urllib.request.Request(url, data=data, method="POST", headers={
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
        "User-Agent": LLM_UA,
        "Accept": "application/json",
    })
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return json.loads(resp.read().decode("utf-8"))


def llm_chat(messages, max_tokens=4096, temperature=0.3,
             groq_model=None, openai_model=None, timeout=90):
    """Chat completion with OpenAI as primary and Groq as automatic fallback.

    Tries providers in order (OpenAI first when its key is set); on a rate-limit
    or transient error it falls through to the next provider, then backs off and
    retries the whole set. Returns the assistant message content (str).
    Raises LLMError only if every provider fails every round.
    """
    providers = []
    if OPENAI_API_KEY:
        providers.append(("openai", OPENAI_URL, OPENAI_API_KEY, openai_model or OPENAI_MODEL))
    if GROQ_API_KEY:
        providers.append(("groq", GROQ_URL, GROQ_API_KEY, groq_model or GROQ_MODEL))
    if not providers:
        raise LLMError(500, "No LLM API key configured. Set OPENAI_API_KEY (preferred) or GROQ_API_KEY in .env.")

    last_err = LLMError(502, "LLM request failed")
    for round_i in range(3):
        for name, url, key, model in providers:
            body = {
                "model": model,
                "max_tokens": max_tokens,
                "temperature": temperature,
                "response_format": {"type": "json_object"},
                "messages": messages,
            }
            try:
                data = _llm_post(url, key, body, timeout=timeout)
                return data["choices"][0]["message"]["content"]
            except urllib.error.HTTPError as e:
                detail = e.read().decode("utf-8", "replace")
                print(f"  LLM {name} HTTP {e.code} (round {round_i + 1}): {detail[:200]}")
                last_err = LLMError(e.code, f"{name} API {e.code}: {detail[:400]}")
            except Exception as e:
                print(f"  LLM {name} error (round {round_i + 1}): {e}")
                last_err = LLMError(502, f"{name}: {e}")
        if round_i < 2:
            time.sleep(min(5 * (2 ** round_i), 20))
    raise last_err


def jira_get(path):
    """Make an authenticated GET to the Jira REST API and return parsed JSON."""
    url = f"{JIRA_BASE_URL}/rest/api/3/{path.lstrip('/')}"
    req = urllib.request.Request(url, headers={
        "Authorization": f"Basic {JIRA_AUTH}",
        "Accept": "application/json",
        "Content-Type": "application/json",
    })
    with urllib.request.urlopen(req, timeout=20) as resp:
        return json.loads(resp.read().decode("utf-8"))


def extract_epic_from_fields(f):
    """Return the epic name (or key) for a Jira issue's fields dict, or '' if none."""
    parent = f.get("parent") or {}
    parent_fields = parent.get("fields") or {}
    parent_type = (parent_fields.get("issuetype") or {}).get("name", "")
    if parent_type == "Epic":
        return parent_fields.get("summary", "")
    epic_key = f.get("customfield_10014") or ""
    return epic_key


_EPIC_KEY_RE = re.compile(r'^[A-Z][A-Z0-9]+-\d+$')


def resolve_epic_name(epic_value):
    """If epic_value looks like a raw issue key (classic/company-managed
    "Epic Link" field), fetch and return the epic's actual name. Otherwise
    return it unchanged — it's already a name (e.g. team-managed projects'
    parent field, resolved directly by extract_epic_from_fields)."""
    if not epic_value or not _EPIC_KEY_RE.match(epic_value):
        return epic_value
    try:
        d = jira_get(f"issue/{epic_value}?fields=summary,customfield_10011")
        ef = d.get("fields", {})
        return ef.get("customfield_10011") or ef.get("summary") or epic_value
    except Exception:
        return epic_value


def jira_issue_to_text(issue, epic_name=None):
    """Convert a Jira issue to plain text suitable for the readiness checker.

    Pass epic_name if it's already been resolved (e.g. batch-resolved
    elsewhere) to avoid a redundant lookup; otherwise it's resolved here.
    """
    fields = issue.get("fields", {})
    key    = issue.get("key", "")
    summary = fields.get("summary", "")
    desc_raw = fields.get("description") or {}

    def extract_text(node):
        if not node:
            return ""
        if isinstance(node, str):
            return node
        t = node.get("type", "")
        text = node.get("text", "")
        children = node.get("content", [])
        parts = [extract_text(c) for c in children]
        joined = "".join(parts) if parts else text
        if t in ("paragraph", "heading"):
            return joined + "\n"
        if t == "listItem":
            return "- " + joined
        if t in ("bulletList", "orderedList"):
            return joined + "\n"
        return joined

    description = extract_text(desc_raw).strip() if isinstance(desc_raw, dict) else str(desc_raw or "").strip()
    issue_type = (fields.get("issuetype") or {}).get("name", "")
    status     = (fields.get("status") or {}).get("name", "")
    priority   = (fields.get("priority") or {}).get("name", "")
    assignee   = ((fields.get("assignee") or {}).get("displayName") or "Unassigned")
    story_points = fields.get("story_points") or fields.get("customfield_10016") or ""
    if epic_name is None:
        epic_name = resolve_epic_name(extract_epic_from_fields(fields))

    lines = [f"[{key}] {summary}"]
    if issue_type: lines.append(f"Type: {issue_type}")
    if status:     lines.append(f"Status: {status}")
    if priority:   lines.append(f"Priority: {priority}")
    if assignee:   lines.append(f"Assignee: {assignee}")
    if story_points: lines.append(f"Story Points: {story_points}")
    if epic_name:  lines.append(f"Parent Epic: {epic_name}")
    if description:
        lines.append("")
        lines.append(description)
    return "\n".join(lines)


def jira_agile_get(path):
    """Authenticated GET to the Jira Agile (Software) REST API."""
    url = f"{JIRA_BASE_URL}/rest/agile/1.0/{path.lstrip('/')}"
    req = urllib.request.Request(url, headers={
        "Authorization": f"Basic {JIRA_AUTH}",
        "Accept": "application/json",
    })
    with urllib.request.urlopen(req, timeout=20) as resp:
        return json.loads(resp.read().decode("utf-8"))


def fetch_issue_comments(key, max_comments=3, max_chars=300):
    """Fetch the most recent comments for a Jira issue. Returns a list of plain-text strings."""
    for attempt in range(3):
        try:
            data = jira_get(f"issue/{key}/comment?maxResults=50&orderBy=-created")
            comments = data.get("comments", [])
            recent = comments[-max_comments:] if len(comments) > max_comments else comments

            def _txt(node):
                if not node: return ""
                if isinstance(node, str): return node
                parts = [_txt(c) for c in node.get("content", [])]
                return (" ".join(p for p in parts if p.strip()) or node.get("text", ""))

            result = []
            for c in recent:
                author = (c.get("author") or {}).get("displayName", "Unknown")
                body_raw = c.get("body") or {}
                text = (_txt(body_raw) if isinstance(body_raw, dict) else str(body_raw)).strip()
                if text:
                    result.append(f"{author}: {text[:max_chars]}")
            return result
        except urllib.error.HTTPError as e:
            if e.code == 429:
                time.sleep(2 ** attempt)  # 1s, 2s, 4s backoff
                continue
            return []
        except Exception:
            return []
    return []


def fetch_board_sprints_and_backlog(bid):
    """Fetch sprints + backlog count for one board. Called per-board, on demand."""
    # Sprints (active + future + last 2 closed)
    try:
        sd = jira_agile_get(f"board/{bid}/sprint?state=active,future&maxResults=10")
        raw_sprints = sd.get("values", [])
    except Exception:
        raw_sprints = []
    try:
        # Jira returns closed sprints oldest-first; fetch last page to get most recent
        probe = jira_agile_get(f"board/{bid}/sprint?state=closed&maxResults=1")
        total_closed = probe.get("total", 0)
        if total_closed > 0:
            start_at = max(0, total_closed - 5)
            cd = jira_agile_get(f"board/{bid}/sprint?state=closed&maxResults=5&startAt={start_at}")
            closed_all = cd.get("values", [])
            closed_all.sort(key=lambda s: s.get("endDate", ""), reverse=True)
            raw_sprints += closed_all[:2]
    except Exception:
        pass

    # Backlog count
    try:
        bl = jira_agile_get(f"board/{bid}/backlog?maxResults=1&fields=summary")
        backlog_count = bl.get("total", 0)
    except Exception:
        backlog_count = 0

    # Per-sprint issue counts — run in parallel
    def sprint_count(sprint_id):
        try:
            d = jira_agile_get(f"sprint/{sprint_id}/issue?maxResults=1&fields=summary")
            return d.get("total", 0)
        except Exception:
            return 0

    sprints = []
    if raw_sprints:
        with ThreadPoolExecutor(max_workers=min(len(raw_sprints), 5)) as ex:
            futures = {ex.submit(sprint_count, s["id"]): s for s in raw_sprints}
            for fut in as_completed(futures):
                s = futures[fut]
                sprints.append({
                    "id": s["id"],
                    "name": s["name"],
                    "state": s.get("state", "future"),
                    "startDate": s.get("startDate", ""),
                    "endDate": s.get("endDate", ""),
                    "goal": s.get("goal", ""),
                    "count": fut.result(),
                })
        order = {"active": 0, "future": 1, "closed": 2}
        sprints.sort(key=lambda x: (order.get(x["state"], 3), x.get("endDate", "") if x["state"] == "closed" else x["name"]), reverse=False)
        # For closed sprints, show most recent first
        active_future = [s for s in sprints if s["state"] != "closed"]
        closed = sorted([s for s in sprints if s["state"] == "closed"], key=lambda x: x.get("endDate", ""), reverse=True)
        sprints = active_future + closed

    return {"sprints": sprints, "backlogCount": backlog_count}


def build_pptx(deck):
    """Build a Blend-branded PPTX from a status deck dict. Returns bytes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io, os

    def rgb(hex6):
        h = hex6.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    C = dict(wb="#053057", nt="#00EDED", lt="#A2F3F3",
             cg="#314550", gr="#1A1A1A", dg="#0B0D0E",
             wh="#FFFFFF", ow="#F4F3F0")

    FONT = "Montserrat"
    ROOT_DIR = os.path.dirname(os.path.abspath(__file__))

    # ── Helpers ──────────────────────────────────────────────────────────────
    def add_rect(slide, x, y, w, h, fill_hex):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
        shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, size=14, color="#FFFFFF",
                 align="left", wrap=True):
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        txb.word_wrap = wrap
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(text)
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        run.font.bold = False
        return txb

    def section_label(slide, text, y=0.45):
        add_text(slide, text, 0.62, y, 11.87, 0.28, size=8, color=C["nt"])
        add_rect(slide, 0.62, y + 0.30, 11.87, 0.016, C["nt"])

    def navy_bg(slide):
        add_rect(slide, 0, 0, 13.33, 7.5, C["wb"])

    def add_logo(slide, x, y, w):
        logo_path = os.path.join(ROOT_DIR, "blend-logo-white.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(x), Inches(y), width=Inches(w))

    def add_footer_text(slide):
        add_text(slide, "Private and Confidential", 9.5, 7.15, 3.7, 0.25,
                 size=7, color=C["lt"], align="right")

    m  = deck.get("metrics") or {}
    pn = deck.get("projectName") or "Project"
    sn = deck.get("sprintName") or ""
    wo = deck.get("weekOf") or ""
    hs = deck.get("healthStatus") or "At Risk"
    sid        = deck.get("sprintId") or ""
    sprint_start_raw = deck.get("sprintStart") or ""
    sprint_end_raw   = deck.get("sprintEnd") or ""

    def _fmt_date(iso):
        if not iso:
            return ""
        try:
            import re
            m2 = re.match(r"(\d{4})-(\d{2})-(\d{2})", iso)
            if m2:
                from datetime import date
                d = date(int(m2.group(1)), int(m2.group(2)), int(m2.group(3)))
                return d.strftime("%b %d, %Y")
        except Exception:
            pass
        return iso[:10]

    sprint_timeline = ""
    if sprint_start_raw or sprint_end_raw:
        parts = [_fmt_date(sprint_start_raw), _fmt_date(sprint_end_raw)]
        sprint_timeline = " – ".join(p for p in parts if p)

    # ── Fresh presentation — no reference PPTX manipulation ──────────────────
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width  = Emu(12192000)   # exact 13.333..." widescreen
    prs.slide_height = Emu(6858000)    # exact 7.5"
    # Remove the built-in type="screen4x3" attribute — mismatches our dims and
    # causes PowerPoint to show a repair prompt on open.
    from pptx.oxml.ns import qn
    sldSz = prs._element.find(qn('p:sldSz'))
    if sldSz is not None:
        sldSz.attrib.pop('type', None)
    blank = prs.slide_layouts[6]   # truly blank layout

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)

    # Gear photo background — pre-crop to 16:9 so it fits the slide exactly
    bg_path = os.path.join(ROOT_DIR, "cover-bg.jpeg")
    if os.path.exists(bg_path):
        import tempfile, os as _os
        from PIL import Image as _PIL
        try:
            with _PIL.open(bg_path) as im:
                iw, ih = im.size
                target_ratio = 12192000 / 6858000  # 16:9 ≈ 1.778
                crop_h = int(iw / target_ratio)
                if crop_h <= ih:
                    y_start = (ih - crop_h) // 2   # centre crop vertically
                    cropped = im.crop((0, y_start, iw, y_start + crop_h))
                else:
                    crop_w = int(ih * target_ratio)
                    x_start = (iw - crop_w) // 2
                    cropped = im.crop((x_start, 0, x_start + crop_w, ih))
                tmp = tempfile.NamedTemporaryFile(suffix=".jpeg", delete=False)
                cropped.save(tmp.name, "JPEG", quality=90)
                tmp.close()
            cover.shapes.add_picture(tmp.name, Emu(0), Emu(0),
                                      width=Emu(12192000), height=Emu(6858000))
            _os.unlink(tmp.name)
        except Exception:
            cover.shapes.add_picture(bg_path, Emu(0), Emu(0),
                                      width=Emu(12192000), height=Emu(6858000))

    # Render oval as PNG using PIL so it clips naturally to slide bounds — no shape overflow
    import tempfile as _tmp, os as _os2
    from PIL import Image as _PILI, ImageDraw as _PILDraw
    W_PX, H_PX = 1920, 1080
    overlay = _PILI.new("RGBA", (W_PX, H_PX), (0, 0, 0, 0))
    draw = _PILDraw.Draw(overlay)
    sx = W_PX / 12192000; sy = H_PX / 6858000
    ox = int(Inches(-1.1) * sx); oy = int(Inches(-0.6) * sy)
    ow = int(Inches(6.6)  * sx); oh = int(Inches(8.7)  * sy)
    draw.ellipse([ox, oy, ox+ow, oy+oh], fill=(5,48,87,255))
    draw.ellipse([ox, oy, ox+ow, oy+oh], outline=(0,237,237,255), width=3)
    _tf = _tmp.NamedTemporaryFile(suffix=".png", delete=False)
    overlay.save(_tf.name, "PNG"); _tf.close()
    cover.shapes.add_picture(_tf.name, Emu(0), Emu(0),
                              width=Emu(12192000), height=Emu(6858000))
    _os2.unlink(_tf.name)

    # Text inside shape
    add_text(cover, pn, 0.75, 1.8, 4.2, 1.4, size=28, color=C["wh"])
    from datetime import date as _date
    created_label = _date.today().strftime("%B %d, %Y")
    add_text(cover, sn, 0.75, 3.25, 4.2, 0.38, size=13, color=C["lt"])
    if sprint_timeline:
        add_text(cover, sprint_timeline, 0.75, 3.65, 4.2, 0.32, size=11, color=C["ow"])
    add_text(cover, "Created On: " + created_label, 0.75, 4.0, 4.2, 0.30, size=10, color="#A2A2A2")
    # Blend360.com footer bottom-left (matching reference)
    add_text(cover, "Blend360.com", 0.5, 7.1, 2.5, 0.28, size=8, color=C["lt"])

    # ── Content slide builder ─────────────────────────────────────────────────
    FOOTER_Y = 6.85

    def new_content_slide():
        sl = prs.slides.add_slide(blank)
        navy_bg(sl)
        add_logo(sl, 0.21, 7.1, 0.9)
        add_footer_text(sl)
        # Thin turquoise top rule
        add_rect(sl, 0, 0, 13.33, 0.04, C["nt"])
        return sl

    blockers     = deck.get("blockers") or []
    has_blockers = len(blockers) > 0
    blocker_h    = 1.0 if has_blockers else 0.0
    col_bottom   = FOOTER_Y - blocker_h

    # ── Slide 2: Exec Summary + Metrics + Accomplishments + Next Steps ────────
    sl = new_content_slide()

    section_label(sl, "EXECUTIVE SUMMARY")
    add_text(sl, deck.get("executiveSummary") or "", 0.62, 0.90, 11.87, 0.70,
             size=12, color=C["ow"])

    stats = [
        ("Total",       str(m.get("totalIssues",      0))),
        ("Completed",   str(m.get("doneIssues",       0))),
        ("In Progress", str(m.get("inProgressIssues", 0))),
        ("Completion",  str(m.get("completionPct",    0)) + "%"),
    ]
    tw, sx, sy_m = 2.83, 0.62, 1.70
    for i, (lbl, val) in enumerate(stats):
        x = sx + i * (tw + 0.15)
        add_rect(sl, x, sy_m, tw, 0.74, C["cg"])
        add_text(sl, val, x, sy_m + 0.04, tw, 0.44, size=24, color=C["nt"], align="center")
        add_text(sl, lbl, x, sy_m + 0.52, tw, 0.20, size=9,  color=C["lt"], align="center")

    col_top = 2.58
    col_h   = col_bottom - col_top
    col_w   = 5.83
    left_x  = 0.62
    right_x = 6.66

    add_text(sl, "KEY ACCOMPLISHMENTS", left_x,  col_top, col_w, 0.26, size=8, color=C["nt"])
    add_rect(sl, left_x,  col_top + 0.26, col_w, 0.015, C["nt"])
    add_text(sl, "NEXT STEPS",          right_x, col_top, col_w, 0.26, size=8, color=C["nt"])
    add_rect(sl, right_x, col_top + 0.26, col_w, 0.015, C["nt"])

    row_h    = 0.44
    max_rows = max(1, int((col_h - 0.3) / row_h))

    for i, item in enumerate((deck.get("accomplishments") or [])[:max_rows]):
        ry = col_top + 0.32 + i * row_h
        add_rect(sl, left_x, ry + 0.06, 0.2, 0.2, C["nt"])
        add_text(sl, "✓", left_x, ry + 0.05, 0.22, 0.22, size=8, color=C["wb"], align="center")
        add_text(sl, item.get("sentence") or item.get("title") or "",
                 left_x + 0.28, ry, col_w - 0.3, 0.4, size=11, color=C["wh"])

    for i, s in enumerate((deck.get("nextSteps") or [])[:max_rows]):
        ry = col_top + 0.32 + i * row_h
        if i % 2 == 0:
            add_rect(sl, right_x, ry, col_w, row_h, C["cg"])
        add_text(sl, "→", right_x + 0.05, ry + 0.06, 0.24, 0.28, size=11, color=C["nt"])
        add_text(sl, s.get("action") or "", right_x + 0.3, ry + 0.04, col_w - 0.35, 0.38,
                 size=11, color=C["wh"])

    if has_blockers:
        bsy = FOOTER_Y - blocker_h
        add_rect(sl, 0, bsy, 13.33, blocker_h, C["gr"])
        add_rect(sl, 0.06, bsy, 0.06, blocker_h, C["cg"])
        add_text(sl, "ISSUES / BLOCKERS", 0.2, bsy + 0.08, 2.3, 0.22, size=7, color=C["lt"])
        i_clr = {"High": "#F87171", "Medium": "#FCD34D", "Low": "#6EE7B7"}
        bw = 10.6 / max(len(blockers[:4]), 1)
        for i, b in enumerate(blockers[:4]):
            bx = 2.6 + i * (bw + 0.1)
            ic = i_clr.get(b.get("impact") or "", "#FCD34D")
            add_rect(sl, bx, bsy + 0.08, bw, 0.85, C["cg"])
            add_rect(sl, bx, bsy + 0.08, 0.05, 0.85, ic)
            add_text(sl, (b.get("impact") or "").upper(), bx + 0.10, bsy + 0.10, 0.8, 0.2,
                     size=7, color=ic)
            add_text(sl, b.get("title") or "", bx + 0.10, bsy + 0.30, bw - 0.15, 0.54,
                     size=9, color=C["ow"])

    # ── Milestones (conditional) ──────────────────────────────────────────────
    milestones = deck.get("milestones") or []
    if milestones:
        sl = new_content_slide()
        section_label(sl, "MILESTONES")
        s_icon = {"complete": "✓", "in_progress": "◎", "upcoming": "○"}
        s_clr  = {"complete": C["nt"], "in_progress": "#FCD34D", "upcoming": C["ow"]}
        for i, ms in enumerate(milestones):
            y  = 1.0 + i * 0.95
            st = ms.get("status") or "upcoming"
            add_rect(sl, 0.62, y + 0.05, 11.87, 0.72, C["cg"])
            add_text(sl, s_icon.get(st, "○"), 0.75, y + 0.08, 0.5, 0.55, size=20,
                     color=s_clr.get(st, C["ow"]))
            add_text(sl, ms.get("name") or "", 1.35, y + 0.18, 9.5, 0.4, size=15, color=C["wh"])
            if ms.get("date"):
                add_text(sl, ms["date"], 11.3, y + 0.18, 1.2, 0.4, size=13,
                         color=C["lt"], align="right")

    # ── Closing slide ─────────────────────────────────────────────────────────
    cl = prs.slides.add_slide(blank)
    navy_bg(cl)

    # Large "Thank you" in turquoise
    add_text(cl, "Thank you", 1.5, 2.6, 10.33, 1.8, size=60, color=C["nt"], align="center")

    # Centered Blend logo
    logo_path = os.path.join(ROOT_DIR, "blend-logo-white.png")
    if os.path.exists(logo_path):
        cl.shapes.add_picture(logo_path, Inches(5.4), Inches(5.0), width=Inches(2.5))

    # ── Save clean ────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class Handler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):
        print("  " + (fmt % args))

    def _send(self, code, body, content_type="application/json"):
        data = body.encode("utf-8") if isinstance(body, str) else body
        self.send_response(code)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate")
        self.end_headers()
        self.wfile.write(data)

    def do_GET(self):
        # ── Jira endpoints ────────────────────────────────────────────────────
        if self.path == "/board.html":
            try:
                with open(os.path.join(ROOT, "board.html"), "rb") as f:
                    self._send(200, f.read(), "text/html; charset=utf-8")
            except FileNotFoundError:
                self._send(404, "board.html not found", "text/plain")
            return

        if self.path == "/api/jira/board-overview":
            # Lightweight board list only — no sprints/backlog here, those are
            # fetched per-board on demand via /api/jira/board-detail so the
            # list itself renders instantly even with 50+ boards.
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"}))
                return
            try:
                raw_boards = []
                start = 0
                while True:
                    page = jira_agile_get(f"board?maxResults=50&startAt={start}")
                    values = page.get("values", [])
                    raw_boards.extend(values)
                    if page.get("isLast", True) or len(values) < 50:
                        break
                    start += 50
                results = []
                for b in raw_boards:
                    loc = b.get("location") or {}
                    project_key = loc.get("projectKey", "")
                    results.append({
                        "id": b["id"],
                        "name": b.get("name", ""),
                        "type": b.get("type", "scrum"),
                        "projectKey": project_key,
                        "projectName": loc.get("projectName", b.get("name", "")),
                        "avatarUrl": loc.get("avatarURI", ""),
                        "url": f"{JIRA_BASE_URL}/jira/software/projects/{project_key}/boards/{b['id']}",
                        "sprints": [],
                        "backlogCount": None,
                    })
                results.sort(key=lambda x: x["projectName"].lower())
                self._send(200, json.dumps(results))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return

        if self.path.startswith("/api/jira/board-detail"):
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"}))
                return
            from urllib.parse import urlparse, parse_qs
            qs = parse_qs(urlparse(self.path).query)
            board_id = qs.get("boardId", [""])[0]
            if not board_id:
                self._send(400, json.dumps({"error": "Need boardId"}))
                return
            try:
                detail = fetch_board_sprints_and_backlog(board_id)
                self._send(200, json.dumps(detail))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return

        if self.path == "/api/jira/status":
            connected = bool(JIRA_AUTH and JIRA_BASE_URL)
            self._send(200, json.dumps({
                "connected": connected,
                "base_url": JIRA_BASE_URL,
            }))
            return

        if self.path.startswith("/api/jira/issues"):
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"}))
                return
            from urllib.parse import urlparse, parse_qs, urlencode
            qs = parse_qs(urlparse(self.path).query)
            search = qs.get("q", [""])[0].strip()
            max_results = int(qs.get("max", ["25"])[0])
            if search:
                jql = f'(summary ~ "{search}" OR text ~ "{search}") ORDER BY updated DESC'
            else:
                jql = "ORDER BY updated DESC"
            try:
                params = urlencode({"jql": jql, "maxResults": max_results,
                                    "fields": "summary,issuetype,status,priority,assignee,customfield_10016,description,parent,customfield_10014"})
                data = jira_get(f"search?{params}")
                issues = []
                for iss in data.get("issues", []):
                    f = iss.get("fields", {})
                    issues.append({
                        "key":   iss["key"],
                        "summary": f.get("summary", ""),
                        "type":  (f.get("issuetype") or {}).get("name", ""),
                        "status": (f.get("status") or {}).get("name", ""),
                        "priority": (f.get("priority") or {}).get("name", ""),
                        "points": f.get("customfield_10016"),
                        "assignee": ((f.get("assignee") or {}).get("displayName") or ""),
                        "epic": extract_epic_from_fields(f),
                    })
                self._send(200, json.dumps({"issues": issues, "total": data.get("total", 0)}))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return

        if self.path.startswith("/api/jira/issue/"):
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"}))
                return
            key = self.path.split("/api/jira/issue/")[-1].split("?")[0].strip()
            try:
                issue = jira_get(f"issue/{key}?fields=summary,issuetype,status,priority,assignee,customfield_10016,description,parent,customfield_10014")
                self._send(200, json.dumps({
                    "key": issue["key"],
                    "text": jira_issue_to_text(issue),
                    "type": (issue["fields"].get("issuetype") or {}).get("name", ""),
                    "epic": resolve_epic_name(extract_epic_from_fields(issue.get("fields", {}))),
                }))
            except urllib.error.HTTPError as e:
                self._send(e.code, json.dumps({"error": f"Jira {e.code}: {e.read().decode('utf-8','replace')}"}))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return
        if self.path.startswith("/api/jira/board-issues"):
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"}))
                return
            from urllib.parse import urlparse, parse_qs
            qs = parse_qs(urlparse(self.path).query)
            board_id  = qs.get("boardId",  [""])[0]
            sprint_id = qs.get("sprintId", [""])[0]
            is_backlog = qs.get("backlog", [""])[0] == "true"
            try:
                fields = "summary,issuetype,status,priority,assignee,customfield_10016,description,parent,customfield_10014"
                if sprint_id:
                    data = jira_agile_get(f"sprint/{sprint_id}/issue?maxResults=50&fields={fields}")
                elif is_backlog and board_id:
                    data = jira_agile_get(f"board/{board_id}/backlog?maxResults=50&fields={fields}")
                else:
                    self._send(400, json.dumps({"error": "Need boardId+backlog=true or sprintId"}))
                    return
                raw_issues = data.get("issues", [])
                issues = []
                for iss in raw_issues:
                    f = iss.get("fields", {})
                    issues.append({
                        "key":      iss["key"],
                        "summary":  f.get("summary", ""),
                        "type":     (f.get("issuetype") or {}).get("name", ""),
                        "status":   (f.get("status")    or {}).get("name", ""),
                        "priority": (f.get("priority")  or {}).get("name", ""),
                        "points":   f.get("customfield_10016"),
                        "assignee": ((f.get("assignee") or {}).get("displayName") or ""),
                        "epic":     extract_epic_from_fields(f),
                    })

                # Batch-resolve epic keys to names for classic projects, deduped
                # so issues sharing an epic don't each trigger their own lookup.
                epic_keys = list(set(
                    i["epic"] for i in issues
                    if i.get("epic") and _EPIC_KEY_RE.match(i["epic"])
                ))
                epic_name_map = {}
                if epic_keys:
                    def fetch_epic_summary(key):
                        try:
                            d = jira_get(f"issue/{key}?fields=summary,customfield_10011")
                            f2 = d.get("fields", {})
                            return key, (f2.get("customfield_10011") or f2.get("summary") or key)
                        except Exception:
                            return key, key
                    with ThreadPoolExecutor(max_workers=min(len(epic_keys), 5)) as pool:
                        for k, name in pool.map(fetch_epic_summary, epic_keys):
                            epic_name_map[k] = name

                # Resolve each issue's epic and bake the resolved name into its
                # text (jira_issue_to_text embeds "Parent Epic: ..." itself, so
                # it needs the resolved name too, not just the "epic" field).
                for entry, iss in zip(issues, raw_issues):
                    resolved = epic_name_map.get(entry["epic"], entry["epic"])
                    entry["epic"] = resolved
                    entry["text"] = jira_issue_to_text(iss, epic_name=resolved)

                self._send(200, json.dumps({"issues": issues, "total": data.get("total", 0)}))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return

        if self.path in ("/status", "/status.html"):
            try:
                with open(os.path.join(ROOT, "status.html"), "rb") as f:
                    self._send(200, f.read(), "text/html; charset=utf-8")
            except FileNotFoundError:
                self._send(404, "status.html not found", "text/plain")
            return

        # ── History endpoints ─────────────────────────────────────────────────
        if self.path == "/api/history":
            self._send(200, json.dumps(load_history()))
            return
        if self.path == "/api/history/clear":
            save_history([])
            self._send(200, json.dumps({"ok": True}))
            return
        if self.path.startswith("/api/history/delete/"):
            entry_id = self.path.split("/")[-1]
            entries = [e for e in load_history() if str(e.get("id")) != entry_id]
            save_history(entries)
            self._send(200, json.dumps({"ok": True}))
            return
        if self.path in ("/", "/landing.html"):
            try:
                with open(os.path.join(ROOT, "landing.html"), "rb") as f:
                    self._send(200, f.read(), "text/html; charset=utf-8")
            except FileNotFoundError:
                self._send(404, "landing.html not found", "text/plain")
            return
        if self.path in ("/checker", "/index.html"):
            try:
                with open(os.path.join(ROOT, "index.html"), "rb") as f:
                    self._send(200, f.read(), "text/html; charset=utf-8")
            except FileNotFoundError:
                self._send(404, "index.html not found", "text/plain")
        elif self.path.startswith("/vendor/"):
            # Serve locally vendored JS libs (no CDN dependency in the browser)
            rel = self.path.lstrip("/").split("?", 1)[0]
            safe = os.path.normpath(rel).replace("\\", "/")
            if not safe.startswith("vendor/"):
                self._send(403, "Forbidden", "text/plain")
                return
            path = os.path.join(ROOT, safe)
            try:
                with open(path, "rb") as f:
                    self._send(200, f.read(), "application/javascript; charset=utf-8")
            except FileNotFoundError:
                self._send(404, "Not found", "text/plain")
        elif any(self.path.endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".jfif", ".webp", ".gif")):
            rel = self.path.lstrip("/").split("?", 1)[0]
            safe = os.path.normpath(rel).replace("\\", "/")
            if "/" in safe:
                self._send(403, "Forbidden", "text/plain")
                return
            path = os.path.join(ROOT, safe)
            try:
                with open(path, "rb") as f:
                    self._send(200, f.read(), "image/png")
            except FileNotFoundError:
                self._send(404, "Not found", "text/plain")
        else:
            self._send(404, "Not found", "text/plain")

    def do_POST(self):
        if self.path == "/api/status-deck":
            if not (OPENAI_API_KEY or GROQ_API_KEY):
                self._send(500, json.dumps({"error": "No LLM API key configured. Set OPENAI_API_KEY or GROQ_API_KEY in .env."})); return
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"})); return
            try:
                length = int(self.headers.get("Content-Length", 0))
                payload = json.loads(self.rfile.read(length) or "{}")
            except Exception as e:
                self._send(400, json.dumps({"error": f"Bad request: {e}"})); return

            board_id    = payload.get("boardId")
            sprint_id   = payload.get("sprintId")
            is_backlog  = payload.get("isBacklog", False)
            project_name = payload.get("projectName", "Project")
            sprint_name  = payload.get("sprintName", "Sprint")
            sprint_start = payload.get("sprintStart", "")
            sprint_end   = payload.get("sprintEnd", "")

            try:
                fields = "summary,issuetype,status,priority,assignee,customfield_10016,labels"
                if sprint_id:
                    data = jira_agile_get(f"sprint/{sprint_id}/issue?maxResults=100&fields={fields}")
                elif is_backlog and board_id:
                    data = jira_agile_get(f"board/{board_id}/backlog?maxResults=50&fields={fields}")
                else:
                    self._send(400, json.dumps({"error": "Need sprintId or boardId+isBacklog=true"})); return

                issues = data.get("issues", [])
                done_issues, in_progress_issues, todo_issues = [], [], []
                sp_planned, sp_done = 0, 0

                for iss in issues:
                    f = iss.get("fields", {})
                    status = f.get("status") or {}
                    cat = (status.get("statusCategory") or {}).get("key", "new")
                    points = f.get("customfield_10016") or 0
                    sp_planned += points or 0
                    desc_raw = f.get("description") or {}
                    def _txt(n):
                        if not n: return ""
                        if isinstance(n, str): return n
                        parts = [_txt(c) for c in n.get("content", [])]
                        return ("; ".join(p for p in parts if p.strip()) if parts else n.get("text", ""))
                    desc_text = _txt(desc_raw).strip()[:300] if isinstance(desc_raw, dict) else str(desc_raw or "").strip()[:300]
                    entry = {
                        "key": iss["key"],
                        "summary": f.get("summary", ""),
                        "type": (f.get("issuetype") or {}).get("name", ""),
                        "status": status.get("name", ""),
                        "priority": (f.get("priority") or {}).get("name", ""),
                        "assignee": ((f.get("assignee") or {}).get("displayName") or "Unassigned"),
                        "points": points,
                        "labels": f.get("labels") or [],
                        "description": desc_text,
                    }
                    if cat == "done":
                        sp_done += points or 0
                        done_issues.append(entry)
                    elif cat == "indeterminate":
                        in_progress_issues.append(entry)
                    else:
                        todo_issues.append(entry)

                total = len(issues)
                done_count = len(done_issues)
                completion_pct = round((done_count / total * 100) if total else 0)

                today = datetime.date.today()
                week_start = today - datetime.timedelta(days=today.weekday())
                week_end   = week_start + datetime.timedelta(days=6)
                week_of    = f"{week_start.strftime('%B %d')}–{week_end.strftime('%d, %Y')}"

                # Cap issues sent to Groq to stay within token limits
                done_capped     = done_issues[:30]
                progress_capped = in_progress_issues[:20]
                todo_capped     = todo_issues[:15]

                # Fetch comments only for in-progress + blocked issues (max 2 each, 120 chars)
                blocked_issues = [i for i in progress_capped + todo_capped
                                  if any(l.lower() in ("blocked", "risk", "blocker")
                                         for l in i.get("labels", []))
                                  or "block" in i.get("status", "").lower()]
                needs_comments = progress_capped + blocked_issues
                comments_map = {}
                if needs_comments:
                    with ThreadPoolExecutor(max_workers=4) as pool:
                        futures = {pool.submit(fetch_issue_comments, iss["key"]): iss["key"]
                                   for iss in needs_comments}
                        for fut in as_completed(futures):
                            key = futures[fut]
                            try:
                                comments_map[key] = [c[:120] for c in fut.result()[:2]]
                            except Exception:
                                comments_map[key] = []

                def fmt(iss, include_desc=False):
                    pts = f"| Points: {iss['points']}" if iss['points'] else ""
                    lbl = f"| Labels: {', '.join(iss['labels'])}" if iss['labels'] else ""
                    desc = ("\n  Description: " + iss['description'][:150] if include_desc and iss.get('description') else "")
                    line = f"- [{iss['key']}] {iss['summary']} | Type: {iss['type']} | Assignee: {iss['assignee']} | Priority: {iss['priority']} {pts} {lbl}{desc}"
                    coms = comments_map.get(iss["key"], [])
                    if coms:
                        line += "\n  Latest comments:\n" + "\n".join(f"    > {c}" for c in coms)
                    return line

                lines = [
                    f"PROJECT: {project_name}", f"SPRINT: {sprint_name}", f"DATE: {week_of}",
                    f"TOTAL: {total} | DONE: {done_count} | IN PROGRESS: {len(in_progress_issues)} | TO DO: {len(todo_issues)}",
                    f"STORY POINTS: Planned={int(sp_planned)} Done={int(sp_done)}",
                    "", "COMPLETED (Done):",
                ] + ([fmt(i, include_desc=True) for i in done_capped] if done_capped else ["(none)"]) + [
                    "", "IN PROGRESS:",
                ] + ([fmt(i) for i in progress_capped] if progress_capped else ["(none)"]) + [
                    "", "TO DO / NOT STARTED:",
                ] + ([fmt(i) for i in todo_capped] if todo_capped else ["(none)"])

                try:
                    text = llm_chat(
                        messages=[
                            {"role": "system", "content": STATUS_DECK_SYSTEM_PROMPT},
                            {"role": "user",   "content": "\n".join(lines)},
                        ],
                        max_tokens=1500,
                        temperature=0.3,
                        groq_model="llama-3.1-8b-instant",
                    )
                except LLMError as e:
                    code = 429 if e.status == 429 else 502
                    self._send(code, json.dumps({
                        "error": "rate_limited" if e.status == 429 else "LLM request failed",
                        "detail": e.detail,
                    }))
                    return

                text = text.replace("```json", "").replace("```", "").strip()
                parsed = json.loads(text)
                parsed["metrics"] = {
                    "totalIssues": total, "doneIssues": done_count,
                    "inProgressIssues": len(in_progress_issues), "toDoIssues": len(todo_issues),
                    "completionPct": completion_pct,
                    "storyPointsPlanned": int(sp_planned), "storyPointsDone": int(sp_done),
                }
                parsed["sprintId"]    = str(sprint_id) if sprint_id else ""
                parsed["sprintStart"] = sprint_start
                parsed["sprintEnd"]   = sprint_end
                self._send(200, json.dumps(parsed))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return

        if self.path == "/api/export-pptx":
            try:
                length = int(self.headers.get("Content-Length", 0))
                deck = json.loads(self.rfile.read(length) or "{}")
            except Exception as e:
                self._send(400, json.dumps({"error": str(e)})); return
            try:
                pptx_bytes = build_pptx(deck)
                safe_name = (deck.get("projectName") or "status").replace(" ", "-")
                fn = safe_name + "-status-deck.pptx"
                self.send_response(200)
                self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                self.send_header("Content-Disposition", f'attachment; filename="{fn}"')
                self.send_header("Content-Length", str(len(pptx_bytes)))
                self.send_header("Access-Control-Allow-Origin", "*")
                self.end_headers()
                self.wfile.write(pptx_bytes)
            except Exception as e:
                import traceback
                self._send(500, json.dumps({"error": str(e), "trace": traceback.format_exc()}))
            return

        if self.path == "/api/jira/update-description":
            if not JIRA_AUTH:
                self._send(503, json.dumps({"error": "Jira not configured"})); return
            try:
                length = int(self.headers.get("Content-Length", 0))
                payload = json.loads(self.rfile.read(length) or "{}")
            except Exception as e:
                self._send(400, json.dumps({"error": str(e)})); return

            key  = (payload.get("key") or "").strip().upper()
            text = (payload.get("description") or "").strip()
            if not key or not text:
                self._send(400, json.dumps({"error": "key and description are required"})); return

            # Convert plain text to Atlassian Document Format (ADF)
            paragraphs = []
            for line in text.split("\n"):
                if line.strip():
                    paragraphs.append({
                        "type": "paragraph",
                        "content": [{"type": "text", "text": line}]
                    })
                else:
                    paragraphs.append({"type": "paragraph", "content": []})
            adf = {"version": 1, "type": "doc", "content": paragraphs or [{"type": "paragraph", "content": []}]}

            body = json.dumps({"fields": {"description": adf}}).encode("utf-8")
            url  = f"{JIRA_BASE_URL}/rest/api/3/issue/{key}"
            req  = urllib.request.Request(url, data=body, method="PUT",
                       headers={"Authorization": f"Basic {JIRA_AUTH}",
                                "Content-Type": "application/json",
                                "Accept": "application/json"})
            try:
                with urllib.request.urlopen(req, timeout=15) as resp:
                    _ = resp.read()
                self._send(200, json.dumps({"ok": True, "key": key}))
            except urllib.error.HTTPError as e:
                detail = e.read().decode("utf-8", "replace")
                self._send(e.code, json.dumps({"error": f"Jira {e.code}", "detail": detail}))
            except Exception as e:
                self._send(502, json.dumps({"error": str(e)}))
            return

        if self.path != "/api/analyze":
            self._send(404, json.dumps({"error": "Not found"}))
            return

        if not (OPENAI_API_KEY or GROQ_API_KEY):
            self._send(500, json.dumps({"error": "No LLM API key configured. Set OPENAI_API_KEY or GROQ_API_KEY in .env."}))
            return

        try:
            length = int(self.headers.get("Content-Length", 0))
            payload = json.loads(self.rfile.read(length) or "{}")
        except Exception as e:
            self._send(400, json.dumps({"error": f"Bad request: {e}"}))
            return

        story = (payload.get("story") or "").strip()
        epic = (payload.get("epic") or "").strip()
        dor = (payload.get("dor") or "").strip()
        mode = (payload.get("mode") or "story").strip()  # "story" | "epic"
        jira_key = (payload.get("jira_key") or "").strip().upper()

        active_prompt = EPIC_SYSTEM_PROMPT if mode == "epic" else SYSTEM_PROMPT

        context_parts = []
        if epic:
            context_parts.append(f"Parent epic: {epic}")
        if dor:
            context_parts.append(f"Team Definition of Ready: {dor}")

        # If this came from Jira, pull recent comments — they often surface
        # clarifying questions, scope disputes, or dependencies the ticket
        # text itself never captured.
        if jira_key and JIRA_AUTH:
            comments = fetch_issue_comments(jira_key, max_comments=5, max_chars=300)
            if comments:
                context_parts.append(
                    "Recent Jira comments (may reveal ambiguity, scope changes, "
                    "or unresolved dependencies not captured in the text above):\n"
                    + "\n".join(f"- {c}" for c in comments)
                )

        # Level 1 + 2: inject team history patterns and similar past stories
        team_ctx = build_team_context(load_history(), mode, story)
        if team_ctx:
            context_parts.append(team_ctx)

        context = "\n\n".join(context_parts)
        user_content = f"User story:\n{story}" + (f"\n\n{context}" if context else "")

        try:
            text = llm_chat(
                messages=[
                    {"role": "system", "content": active_prompt},
                    {"role": "user",   "content": user_content},
                ],
                max_tokens=4096,
                temperature=0.3,
                groq_model=GROQ_MODEL,
            )
        except LLMError as e:
            code = 429 if e.status == 429 else 502
            self._send(code, json.dumps({
                "error": "rate_limited" if e.status == 429 else "LLM request failed",
                "detail": e.detail,
            }))
            return

        try:
            text = text.replace("```json", "").replace("```", "").strip()
            parsed = json.loads(text)
        except Exception as e:
            self._send(502, json.dumps({"error": f"Could not parse model output: {e}", "raw": text}))
            return

        # Enforce consistent readiness_level based on total score
        total = parsed.get("total", 0)
        if total >= 80:
            parsed["readiness_level"] = "Decomposition Ready" if mode == "epic" else "Sprint Ready"
        elif total >= 60:
            parsed["readiness_level"] = "Almost Ready"
        elif total >= 40:
            parsed["readiness_level"] = "Needs Work"
        else:
            parsed["readiness_level"] = "Not Ready"

        # Save to history (full result stored for panel view)
        entry = {
            "id": int(time.time() * 1000),
            "title": story.split("\n")[0][:80],
            "story_text": story,
            "mode": mode,
            "epic": epic,
            "score": parsed.get("total", 0),
            "readiness_level": parsed.get("readiness_level", ""),
            "checked_at": datetime.datetime.now(EASTERN_TZ).strftime("%Y-%m-%d %H:%M"),
            "result": parsed,
        }
        if jira_key:
            entry["jira_key"] = jira_key
        history = load_history()
        history.insert(0, entry)
        save_history(history)

        self._send(200, json.dumps(parsed))


if __name__ == "__main__":
    port = int(os.environ.get("PORT", "8000"))
    if not (OPENAI_API_KEY or GROQ_API_KEY):
        print("WARNING: no LLM API key found in .env (set OPENAI_API_KEY and/or GROQ_API_KEY).")
    print(f"LLM -> OpenAI (primary): {'configured' if OPENAI_API_KEY else 'MISSING'} | Groq (fallback): {'configured' if GROQ_API_KEY else 'MISSING'}")
    host = os.environ.get("HOST", "127.0.0.1")
    print(f"Serving on http://{host}:{port}  (Ctrl+C to stop)")
    ThreadingHTTPServer((host, port), Handler).serve_forever()
